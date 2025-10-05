import os
import streamlit as st
import pandas as pd
import numpy as np
import json
from datetime import datetime, timedelta
import plotly.express as px
import requests

try:
    import gspread
    from oauth2client.service_account import ServiceAccountCredentials
    GSHEETS_AVAILABLE = True
except Exception:
    GSHEETS_AVAILABLE = False


DATA_DIR = "./gf_data"
LOCAL_LEADERBOARD = os.path.join(DATA_DIR, "leaderboard.json")
PITCH_DIR = os.path.join(DATA_DIR, "pitch_slides")
os.makedirs(DATA_DIR, exist_ok=True)
os.makedirs(PITCH_DIR, exist_ok=True)

CROPS = {
    "Wheat": {"base_yield": 1.0, "water_need": 1.0, "sensitivity_temp": 0.8},
    "Corn": {"base_yield": 1.2, "water_need": 1.3, "sensitivity_temp": 1.0},
    "Sunflower": {"base_yield": 0.9, "water_need": 0.9, "sensitivity_temp": 0.7},
    "Carrot": {"base_yield": 0.8, "water_need": 0.7, "sensitivity_temp": 0.6},
}

LOCATIONS = {
    "Andijan, Uzbekistan": {"lat":40.7833, "lon":72.3500},
    "Tashkent, Uzbekistan": {"lat":41.2995, "lon":69.2401},
    "Nairobi, Kenya": {"lat":-1.2921, "lon":36.8219},
    "Texas, USA": {"lat":31.9686, "lon":-99.9018},
}

POWER_URL = "https://power.larc.nasa.gov/api/temporal/daily/point"


def load_local_lb():
    if os.path.exists(LOCAL_LEADERBOARD):
        try:
            with open(LOCAL_LEADERBOARD, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            return []
    return []


def save_local_lb(data):
    with open(LOCAL_LEADERBOARD, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def add_to_leaderboard(entry: dict):
    
    gs_path = os.getenv('GSHEET_CREDENTIALS_PATH')
    sheet_name = os.getenv('GSHEET_NAME', 'GF_Leaderboard')
    if gs_path and GSHEETS_AVAILABLE:
        try:
            scope = ['https://spreadsheets.google.com/feeds', 'https://www.googleapis.com/auth/drive']
            creds = ServiceAccountCredentials.from_json_keyfile_name(gs_path, scope)
            client = gspread.authorize(creds)
            sh = client.open(sheet_name)
            worksheet = sh.sheet1
            
            row = [entry.get('date',''), entry.get('player',''), entry.get('location',''), entry.get('crop',''), entry.get('score',0), entry.get('revenue',0), entry.get('profit',0), entry.get('sustainability',0)]
            worksheet.append_row(row)
            return True
        except Exception as e:
            st.warning(f"Google Sheets leaderboard failed, falling back to local. ({e})")
    
    data = load_local_lb()
    data.append(entry)
    data = sorted(data, key=lambda x: x.get('score',0), reverse=True)[:200]
    save_local_lb(data)
    return False

def fetch_power_daily(lat, lon, start, end):
    """Fetch daily POWER data for given coordinates and date range.
    Returns DataFrame with date, T2M (temp), PRECTOTCORR (precip), ALLSKY_SFC_SW_DWN (solar)
    """
    params = {
        'start': start.strftime('%Y%m%d'),
        'end': end.strftime('%Y%m%d'),
        'latitude': lat,
        'longitude': lon,
        'community': 'ag',
        'parameters': 'T2M,PRECTOTCORR,ALLSKY_SFC_SW_DWN',
        'format': 'JSON'
    }
    try:
        r = requests.get(POWER_URL, params=params, timeout=20)
        r.raise_for_status()
        data = r.json()
        days = []
        values = data['properties']['parameter']
        dates = sorted(values['T2M'].keys())
        for d in dates:
            days.append({
                'date': pd.to_datetime(d),
                'temp': float(values['T2M'][d]),
                'rainfall': float(values['PRECTOTCORR'][d]) if 'PRECTOTCORR' in values else 0.0,
                'solar': float(values['ALLSKY_SFC_SW_DWN'][d]) if 'ALLSKY_SFC_SW_DWN' in values else 0.0
            })
        return pd.DataFrame(days)
    except Exception as e:
        st.warning(f"POWER API fetch failed: {e} — using mock data.")
        return None


def generate_mock_nasa_data(location_name: str, start_date: str = None, days=180):
    if start_date is None:
        start = pd.to_datetime(datetime.utcnow().date())
    else:
        start = pd.to_datetime(start_date)
    np.random.seed(abs(hash(location_name)) % (2**32))
    dates = [start + timedelta(days=i) for i in range(days)]
    base_temp = 20 + 8 * np.sin(np.linspace(0, 2*np.pi, days))
    temp = base_temp + np.random.normal(0, 2, days)
    rainfall = np.clip(2 + 8 * np.sin(np.linspace(1, 3*np.pi, days)) + np.random.normal(0,5,days), 0, None)
    soil = np.zeros(days)
    soil[0] = 0.3 + 0.1*np.random.rand()
    for i in range(1, days):
        soil[i] = soil[i-1] + 0.01*rainfall[i] - 0.004*max(0, temp[i]-20)
        soil[i] = np.clip(soil[i], 0.05, 0.9)
    ndvi = np.clip(0.2 + 0.6 * soil + np.random.normal(0,0.05,days), 0, 1)
    df = pd.DataFrame({'date': dates, 'rainfall': rainfall, 'temp': temp, 'soil_moisture': soil, 'ndvi': ndvi})
    return df

def simulate_day(state: dict, action: dict, obs: dict):
    crop = state['crop']
    meta = CROPS[crop]
    soil = obs.get('soil_moisture', 0.3)
    temp = obs.get('temp', 20)
    rain = obs.get('rainfall', 0)
    ndvi = obs.get('ndvi', 0.4)
    water_applied = action.get('water', 1)
    water_score = (soil + 0.15*rain/10.0) * (1 + 0.15*water_applied)
    temp_score = max(0.4, 1 - abs(temp - 22)/25 * meta['sensitivity_temp'])
    fert = action.get('fertilize', 0)
    fert_score = 1.0 + 0.12*fert
    pest = action.get('pesticide', 0)
    pest_score = 1.0 + 0.08*pest
    health = state.get('health', 0.6)
    delta = 0.025 * (water_score * temp_score * fert_score * pest_score - 1)
    health = np.clip(health + delta, 0.01, 1.0)
    water_cost = 0.12 * state['area'] * water_applied
    fert_cost = 0.55 * state['area'] * fert
    pest_cost = 0.35 * state['area'] * pest
    state['health'] = float(health)
    state['money'] = float(state['money'] - water_cost - fert_cost - pest_cost)
    state['day_index'] = state.get('day_index',0) + 1
    state['total_water'] = state.get('total_water',0) + water_applied*state['area']
    state['total_fert'] = state.get('total_fert',0) + fert*state['area']
    state['total_pest'] = state.get('total_pest',0) + pest*state['area']
    metrics = {'health':health,'ndvi':ndvi,'soil':soil,'temp':temp,'rain':rain,'water_cost':water_cost,'fert_cost':fert_cost,'pest_cost':pest_cost}
    return state, metrics


def finalize_harvest(state: dict):
    meta = CROPS[state['crop']]
    yield_factor = state.get('health',0.5)
    gross_yield = meta['base_yield'] * state['area'] * yield_factor * (1 + 0.05*np.random.randn())
    price_per_unit = 10.0 * meta['base_yield']
    revenue = max(0, gross_yield * price_per_unit)
    spent = state['initial_money'] - state['money']
    profit = revenue - spent
    water_used = state.get('total_water',0)
    fert_used = state.get('total_fert',0)
    pest_used = state.get('total_pest',0)
    sustainability = max(0, 1 - 0.02*water_used - 0.05*fert_used - 0.03*pest_used)
    score = max(0, revenue * 0.5 + sustainability * 120 + max(0, profit) * 0.2)
    badges = []
    if revenue > 50*state['area']:
        badges.append('Top Producer')
    if sustainability > 0.7:
        badges.append('Nature Protector')
    return {'gross_yield':float(gross_yield),'revenue':float(revenue),'profit':float(profit),'sustainability':float(sustainability),'score':float(score),'badges':badges}

def generate_pitch_pptx(player, location, crop):
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "GlobalFarming — Pitch"
        subtitle = slide.placeholders[1]
        subtitle.text = f"{player} — {location} — {crop}"
        
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Problem & Opportunity"
        s2.placeholders[1].text = "Smallholder farmers lack accessible tools to use satellite and climate data. NASA datasets can bridge this gap."
        
        s3 = prs.slides.add_slide(prs.slide_layouts[1])
        s3.shapes.title.text = "Our Solution"
        s3.placeholders[1].text = "An educational game that uses NASA data to train farmers in sustainable practices, with leaderboards and real-time dashboards."
        filename = os.path.join(PITCH_DIR, f"pitch_{player}_{int(datetime.utcnow().timestamp())}.pptx")
        prs.save(filename)
        return filename
    except Exception as e:
        return None

def local_css():
    st.markdown("""
    <style>
    .big-header{font-size:34px;color:#0b3d0b;font-weight:700}
    .sub-header{font-size:16px;color:#2e7d32}
    .card{background:linear-gradient(180deg,#ffffff,#f1f8e9);padding:16px;border-radius:12px;box-shadow:0 4px 12px rgba(0,0,0,0.08)}
    .footer{color:#555;font-size:13px}
    </style>
    """, unsafe_allow_html=True)


def main():
    st.set_page_config(page_title="GlobalFarming — For Modern Farmers", layout='wide', page_icon='🌾')
    local_css()
    if 'state' not in st.session_state:
        st.session_state['state'] = 'landing'
    if st.session_state['state'] == 'landing':
        landing()
    elif st.session_state['state'] == 'menu':
        menu()
    elif st.session_state['state'] == 'instructions':
        instructions()
    elif st.session_state['state'] == 'info':
        info_dashboard()
    elif st.session_state['state'] == 'setup':
        setup_game()
    elif st.session_state['state'] == 'play':
        play()
    elif st.session_state['state'] == 'results':
        results()


def landing():
    st.markdown('<div class="big-header">🌍 GlobalFarming — For Modern Farmers</div>', unsafe_allow_html=True)
    st.markdown('<p class="sub-header">Data-driven educational farming game using NASA datasets — learn, play, and protect nature.</p>', unsafe_allow_html=True)
    col1, col2 = st.columns([2,1])
    with col1:
        st.image('https://images.unsplash.com/photo-1501004318641-b39e6451bec6?auto=format&fit=crop&w=1200&q=60')
    with col2:
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('**Welcome!**')
        st.write('Start your farming adventure, explore past-season dashboards powered by NASA data, or read the game instructions.')
        if st.button('Start'):
            st.session_state['state'] = 'menu'
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)
    st.markdown('<hr>')
    st.write('Tip: To enable cloud leaderboard, set environment variable GSHEET_CREDENTIALS_PATH pointing to your Google service account JSON file and set GSHEET_NAME.')


def menu():
    st.header('Main Menu')
    c1, c2, c3 = st.columns(3)
    with c1:
        if st.button('1) Instructions'):
            st.session_state['state'] = 'instructions'
            st.rerun()
    with c2:
        if st.button('2) Information (Dashboard)'):
            st.session_state['state'] = 'info'
            st.rerun()
    with c3:
        if st.button('3) Start Game'):
            st.session_state['state'] = 'setup'
            st.rerun()
    st.subheader('Leaderboard Preview')
    lb = load_local_lb()
    if lb:
        df = pd.DataFrame(lb).head(10)
        st.table(df[['player','location','crop','score']])
    else:
        st.info('No leaderboard entries yet — be the first!')


def instructions():
    st.header('How to Play — Quick Guide')
    st.markdown('''
    - **Goal:** Maximize yield & profit **and** minimize environmental impact.
    - **Flow:** Choose location & crop → daily actions (water/fertilize/pesticide) → use NASA-powered dashboard for insights → harvest → leaderboard.
    - **Scoring:** Composite of Yield, Profit, Sustainability. Badges for Top Producer and Nature Protector.
    ''')
    st.markdown('---')
    if st.button('Back to Menu'):
        st.session_state['state'] = 'menu'
        st.rerun()


def info_dashboard():
    st.header('Information — Past Seasons Dashboard')
    col1, col2 = st.columns([1,2])
    with col1:
        loc = st.selectbox('Choose location', list(LOCATIONS.keys()))
        start = st.date_input('Start date', value=datetime(2024,1,1))
        days = st.number_input('Days (history length)', min_value=30, max_value=365, value=180)
        if st.button('Load Data'):
            coords = LOCATIONS[loc]
            df = fetch_power_daily(coords['lat'], coords['lon'], pd.to_datetime(start), pd.to_datetime(start)+timedelta(days=days-1))
            if df is None:
                df = generate_mock_nasa_data(loc, start_date=str(start), days=days)
           
            if 'soil_moisture' not in df.columns:
                df['soil_moisture'] = np.clip(0.3 + 0.01*np.cumsum(df['rainfall'] - 2), 0.05, 0.9)
            st.session_state['dashboard_df'] = df
            st.session_state['dashboard_loc'] = loc
            st.rerun()
    with col2:
        if 'dashboard_df' in st.session_state:
            df = st.session_state['dashboard_df']
            st.subheader(f"Location: {st.session_state.get('dashboard_loc')}")
            fig = px.line(df, x='date', y=['ndvi' if 'ndvi' in df.columns else 'soil_moisture','soil_moisture'], labels={'value':'metric'})
            st.plotly_chart(fig, use_container_width=True)
            st.write('Summary stats:')
            st.metric('Avg Soil Moisture', f"{df['soil_moisture'].mean():.2f}")
            st.metric('Total Rainfall (mm)', f"{df['rainfall'].sum():.0f}")
            
            crop = st.selectbox('Choose crop (proxy estimates)', list(CROPS.keys()))
            proxy = CROPS[crop]['base_yield'] * df['soil_moisture'].mean() * 10
            st.info(f'Estimated historical yield index for {crop}: {proxy:.2f}')
        else:
            st.info('Load data from the left to view dashboard')
    st.markdown('---')
    if st.button('Back to Menu'):
        st.session_state['state'] = 'menu'
        st.rerun()

def get_nasa_data(lat, lon, date=None):
    """
    Fetch real NASA POWER data for the given location and date.
    Cleans invalid (-999, -9999) values automatically.
    """
    if date is None:
        date = datetime.today().strftime("%Y%m%d")

    base_url = "https://power.larc.nasa.gov/api/temporal/daily/point"

    params = {
        "parameters": "T2M,TS,WS10M,QV2M",  
        "community": "AG",
        "longitude": lon,
        "latitude": lat,
        "start": date,
        "end": date,
        "format": "JSON"
    }

    try:
        response = requests.get(base_url, params=params)
        response.raise_for_status()
        data = response.json()

        t2m = data["properties"]["parameter"]["T2M"][date]
        ts = data["properties"]["parameter"]["TS"][date]
        ws = data["properties"]["parameter"]["WS10M"][date]
        qv = data["properties"]["parameter"]["QV2M"][date]

        def clean(val):
            return None if val in [-999, -9999, -999.0, -9999.0] else round(val, 2)

        t2m = clean(t2m)
        ts = clean(ts)
        ws = clean(ws)
        qv = clean(qv)

        
        soil_moisture = round(0.1 + (qv or 0) * 0.01, 2)
        ndvi = round(0.3 + (ws or 0) * 0.02, 2)

        return {
            "temperature": t2m,
            "surface_temp": ts,
            "wind_speed": ws,
            "humidity": qv,
            "soil_moisture": soil_moisture,
            "ndvi": ndvi,
        }

    except Exception as e:
        print("NASA API Error:", e)
        return {
            "temperature": None,
            "surface_temp": None,
            "wind_speed": None,
            "humidity": None,
            "soil_moisture": None,
            "ndvi": None,
        }


def setup_game():
    st.header('Game Setup')
    name = st.text_input('Player name', value=f'Farmer_{np.random.randint(100,999)}')
    location = st.selectbox('Location', list(LOCATIONS.keys()))
    crop = st.selectbox('Crop', list(CROPS.keys()))
    area = st.number_input('Area (ha)', min_value=0.1, max_value=100.0, value=1.0)
    start_money = st.number_input('Starting capital', value=200.0)
    seed_info = st.checkbox('Show data-driven tips before starting')
    if seed_info and 'dashboard_df' in st.session_state:
        df = st.session_state['dashboard_df']
        st.write('Pre-game insights based on loaded dashboard:')
        st.metric('Avg Soil Moisture', f"{df['soil_moisture'].mean():.2f}")
    if st.button('Begin Adventure'):
        st.session_state['game'] = {'player':name,'location':location,'crop':crop,'area':float(area),'money':float(start_money),'initial_money':float(start_money),'day_index':0,'health':0.6,'history':[],'total_water':0,'total_fert':0,'total_pest':0}
       
        coords = LOCATIONS[location]
        df = fetch_power_daily(coords['lat'], coords['lon'], datetime.utcnow().date(), datetime.utcnow().date()+timedelta(days=179))
        if df is None or df.empty:
            df = generate_mock_nasa_data(location, start_date=str(datetime.utcnow().date()), days=180)
        
        if 'soil_moisture' not in df.columns:
            df['soil_moisture'] = np.clip(0.3 + 0.01*np.cumsum(df['rainfall'] - 2), 0.05, 0.9)
        if 'ndvi' not in df.columns:
            df['ndvi'] = np.clip(0.2 + 0.6*df['soil_moisture'] + np.random.normal(0,0.05,len(df)), 0,1)
        st.session_state['nasa_df'] = df
        st.session_state['state'] = 'play'
        st.rerun()
    if st.button('Back to Menu'):
        st.session_state['state'] = 'menu'
        st.rerun()



def play():
    game = st.session_state['game']
    st.header(f"🌾 Farming — {game['player']} | {game['location']} | {game['crop']}")

    left, right = st.columns([1, 2])

    with left:
        st.subheader('🎮 Today actions')
        day_idx = game['day_index']

        coords = LOCATIONS[game['location']]
        lat, lon = coords['lat'], coords['lon']

        
        nasa_data = get_nasa_data(lat, lon)

        sim_date = datetime.utcnow().date()

        st.write("Simulated date:", f"📅 {sim_date}")

        sm = nasa_data["soil_moisture"]
        ndvi = nasa_data["ndvi"]
        temp = nasa_data["temperature"]

        st.metric("Soil Moisture", f"{sm:.2f}" if sm else "No Data")
        st.metric("NDVI", f"{ndvi:.2f}" if ndvi else "No Data")
        st.metric("Temp (°C)", f"{temp:.1f}" if temp else "No Data 🌥️")

        water = st.radio('💧 Watering level', ['Low', 'Medium', 'High'], index=1)
        water_map = {'Low': 0, 'Medium': 1, 'High': 2}
        fert = st.checkbox('🌱 Fertilize')
        pest = st.checkbox('🪲 Apply pesticide')

        if st.button('✅ Execute'):
            action = {
                'water': water_map[water],
                'fertilize': int(fert),
                'pesticide': int(pest)
            }
            game['total_water'] += action['water'] * game['area']
            game['total_fert'] += action['fertilize'] * game['area']
            game['total_pest'] += action['pesticide'] * game['area']

            game, metrics = simulate_day(
                game, action, {
                    'soil_moisture': sm or 0.2,
                    'temp': temp or 20,
                    'rainfall': 0.1,
                    'ndvi': ndvi or 0.3
                }
            )

            game['history'].append({
                'date': str(sim_date),
                'action': action,
                'metrics': metrics,
                'money': game['money']
            })

            st.session_state['game'] = game
            st.success('✅ Actions executed for the day!')
            st.rerun()

        if st.button('🌾 Harvest & Finish'):
            results = finalize_harvest(game)
            entry = {
                'player': game['player'],
                'location': game['location'],
                'crop': game['crop'],
                'date': datetime.utcnow().isoformat(),
                'score': results['score'],
                'revenue': results['revenue'],
                'profit': results['profit'],
                'sustainability': results['sustainability']
            }
            add_to_leaderboard(entry)
            st.session_state['results'] = results
            st.session_state['state'] = 'results'
            st.rerun()

        if st.button('🚪 Abort'):
            st.session_state['state'] = 'menu'
            st.rerun()

    with right:
        st.subheader('🗺️ Farm Overview')
        coords = LOCATIONS[game['location']]
        st.map(pd.DataFrame({'lat': [coords['lat']], 'lon': [coords['lon']]}))

        if game['history']:
            hist = pd.DataFrame([
                {'date': h['date'], 'health': h['metrics']['health'], 'money': h['money']}
                for h in game['history']
            ])
            st.plotly_chart(px.line(hist, x='date', y=['health', 'money'], title='Farm Progress'), use_container_width=True)
            st.dataframe(hist)
        else:
            st.info('📈 No actions yet — execute daily actions to grow your farm!')

def results():
    st.header('Harvest Results')
    r = st.session_state['results']
    st.metric('Gross Yield', f"{r['gross_yield']:.2f}")
    st.metric('Revenue', f"${r['revenue']:.2f}")
    st.metric('Profit', f"${r['profit']:.2f}")
    st.metric('Sustainability', f"{r['sustainability']:.2f}")
    st.metric('Final Score', f"{r['score']:.2f}")
    if r.get('badges'):
        st.success('Badges: ' + ', '.join(r['badges']))
    st.markdown('---')
    st.subheader('Leaderboard (Top Local)')
    lb = load_local_lb()
    if lb:
        st.dataframe(pd.DataFrame(lb).head(20))
    else:
        st.write('No entries yet')
    if st.button('Generate Pitch (PPTX)'):
        g = st.session_state['game']
        fname = generate_pitch_pptx(g['player'], g['location'], g['crop'])
        if fname:
            st.success(f'Pitch saved to {fname}')
        else:
            st.warning('pptx generation not available (python-pptx not installed)')
    if st.button('Back to Menu'):
        st.session_state['state'] = 'menu'
        st.rerun()

if __name__ == '__main__':
    main()
